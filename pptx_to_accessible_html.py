#!/usr/bin/env python3
"""
pptx_to_accessible_html.py
Converts PowerPoint (.pptx) files to accessible HTML, preserving:
  - Heading structure (slide titles → <h1>, content titles → <h2>/<h3>)
  - All text content with proper semantics
  - Images with SME-authored alt text from PowerPoint
  - Slide speaker notes
  - Tables with proper <th>/<td> structure

Usage:
  Single file:      python pptx_to_accessible_html.py presentation.pptx
  Output name:      python pptx_to_accessible_html.py presentation.pptx -o output.html
  Batch folder:     python pptx_to_accessible_html.py ./slides_folder/
  With notes:       python pptx_to_accessible_html.py presentation.pptx --include-notes
  Focusable math:   python pptx_to_accessible_html.py presentation.pptx --focusable-math
  With MathJax:     python pptx_to_accessible_html.py presentation.pptx --mathjax
  Scale images:     python pptx_to_accessible_html.py presentation.pptx --img-scale 60

Requirements:
  pip install python-pptx
"""

import argparse
import base64
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Install it with:  pip install python-pptx")
    sys.exit(1)


# ── Helpers ──────────────────────────────────────────────────────────────────

def get_alt_text(shape) -> str:
    """Extract SME-authored alt text from a shape's XML (the 'descr' attribute)."""
    try:
        # Alt text lives in <p:cNvPr descr="..."> inside nvPicPr or nvSpPr
        nvPr = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr")
        # Use the standard DrawingML namespace path
        cNvPr = shape._element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        )
        # Fallback: walk up through nvPicPr / nvSpPr
        for tag in [
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr",
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr",
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr",
        ]:
            container = shape._element.find(f".//{tag}")
            if container is not None:
                cNvPr = container.find(
                    "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
                )
                if cNvPr is not None:
                    break

        # Most reliable: search all cNvPr elements under the shape
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for cNvPr in shape._element.iter(f"{{{ns}}}cNvPr"):
            descr = cNvPr.get("descr", "").strip()
            if descr:
                return descr

        # Also check the drawing namespace used for grouped/linked images
        ns2 = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for cNvPr in shape._element.iter(f"{{{ns2}}}cNvPr"):
            descr = cNvPr.get("descr", "").strip()
            if descr:
                return descr

        return ""
    except Exception:
        return ""


def get_alt_text_and_decorative(shape) -> tuple:
    """
    Returns (alt_text, is_decorative).

    PowerPoint 365 stores the decorative flag as an XML extension element:
      <adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2016/decorative" val="1"/>
    inside <a:extLst> under <p:cNvPr>. Older versions (and LibreOffice) may
    instead set decorative="1" directly as an attribute on <p:cNvPr>.
    We check both.

    Logic:
      - decorative flag found (either format) → is_decorative=True, alt=""
      - descr attribute is non-empty → is_decorative=False, alt=descr
      - neither → is_decorative=False, alt="" (caller falls back to generated text)
    """
    DECORATIVE_NS = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    try:
        # ── Method 1: PowerPoint 365 extLst extension ──
        for dec_el in shape._element.iter(f"{{{DECORATIVE_NS}}}decorative"):
            if dec_el.get("val", "0") == "1":
                return ("", True)

        # ── Method 2: direct attribute on cNvPr (older PowerPoint / LibreOffice) ──
        for el in shape._element.iter():
            if el.tag.endswith("}cNvPr") or el.tag == "cNvPr":
                if el.get("decorative", "0") == "1":
                    return ("", True)
                descr = el.get("descr", "").strip()
                if descr:
                    return (descr, False)

        return ("", False)
    except Exception:
        return ("", False)


def get_alt_text_reliable(shape) -> str:
    """
    More robust alt-text extractor that searches ALL cNvPr elements
    in the shape's XML tree regardless of namespace prefix used by the file.
    """
    try:
        # Iterate every element in the shape's XML subtree
        for el in shape._element.iter():
            # Match any cNvPr tag regardless of namespace
            if el.tag.endswith("}cNvPr") or el.tag == "cNvPr":
                descr = el.get("descr", "").strip()
                if descr:
                    return descr
        return ""
    except Exception:
        return ""


def image_to_data_uri(image_blob: bytes, content_type: str) -> str:
    """Convert raw image bytes to a base64 data URI for embedding in HTML."""
    b64 = base64.b64encode(image_blob).decode("ascii")
    return f"data:{content_type};base64,{b64}"


def escape_html(text: str) -> str:
    """Minimal HTML escaping."""
    return (
        text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
    )


MONOSPACE_FONTS = {
    "consolas", "courier", "courier new", "lucida console", "monaco",
    "menlo", "source code pro", "fira code", "fira mono",
    "dejavu sans mono", "liberation mono", "roboto mono",
    "jetbrains mono", "inconsolata", "cascadia code", "cascadia mono",
    "sf mono", "hack", "droid sans mono", "ubuntu mono",
    "noto sans mono", "ibm plex mono", "anonymous pro",
    "ocr a", "ocr b", "andale mono", "lucida sans typewriter",
}


def is_code_shape(shape) -> bool:
    """Detect code content via shape name convention or monospace font usage."""
    try:
        if shape.name and "code" in shape.name.lower():
            return True
    except Exception:
        pass
    if not shape.has_text_frame:
        return False
    mono_runs = 0
    total_runs = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            total_runs += 1
            font_name = run.font.name
            if font_name and font_name.strip().lower() in MONOSPACE_FONTS:
                mono_runs += 1
    return total_runs > 0 and mono_runs == total_runs


def render_code_block(tf) -> str:
    """Render a text frame as a <pre><code> block, preserving whitespace."""
    lines = []
    for para in tf.paragraphs:
        # Preserve original text including leading whitespace
        lines.append(escape_html(para.text))
    code_text = "\n".join(lines)
    return f"<pre><code>{code_text}</code></pre>"


def is_title_placeholder(shape) -> bool:
    """Return True if the shape is a title or centered-title placeholder."""
    from pptx.enum.text import PP_ALIGN
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        # Placeholder type 1 = TITLE, 3 = CENTER_TITLE, 15 = SUBTITLE
        return ph.type in (1, 3)
    except Exception:
        return False


def is_subtitle_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type == 15  # SUBTITLE only (not BODY=2, which is main content)
    except Exception:
        return False


def is_slide_number_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type == 13  # SLIDE_NUMBER
    except Exception:
        return False


def is_running_header(shape, deck_title: str) -> bool:
    """Return True if shape is a body placeholder whose text matches the deck title."""
    if not deck_title:
        return False
    try:
        ph = shape.placeholder_format
        if ph is None or ph.type != 2:  # BODY
            return False
        if shape.has_text_frame:
            return shape.text_frame.text.strip() == deck_title
    except Exception:
        pass
    return False


def shape_has_text(shape) -> bool:
    return shape.has_text_frame and shape.text_frame.text.strip()


_CITATION_RE = re.compile(r'^[\d,\s\-–]+$')


# ── OMML → MathML conversion ──────────────────────────────────────────────────
#
# PowerPoint embeds maths as OMML (Office Math Markup Language) inside
# <a14:m><m:oMath>…</m:oMath></a14:m> wrappers that appear as siblings of
# <a:r> text runs inside <a:p> paragraphs.  The entire shape containing math
# is often wrapped in <mc:AlternateContent> and therefore invisible to the
# normal python-pptx shape iterator; we handle that separately in convert_slide.
#
# Namespace URIs
_DML_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
_M_NS    = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A14_NS  = "http://schemas.microsoft.com/office/drawing/2010/main"
_MC_NS   = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_PML_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"

_A   = f"{{{_DML_NS}}}"
_M   = f"{{{_M_NS}}}"
_A14 = f"{{{_A14_NS}}}"
_MC  = f"{{{_MC_NS}}}"
_P   = f"{{{_PML_NS}}}"

# Single-character strings we classify as math operators
_MATH_OPS = frozenset(
    '+-=<>≤≥≠±∓·∘∈∉⊆⊇⊂⊃∪∩∧∨¬→⟹←↔⟺∀∃∂∇∑∏∫∮√∞×÷|/\\:;!?,.^~'
)

# OMML property / control elements whose subtrees carry no renderable content
_M_PROP_TAGS = {
    f"{_M}rPr", f"{_M}sSubSupPr", f"{_M}sSubPr", f"{_M}sSupPr",
    f"{_M}fPr", f"{_M}radPr", f"{_M}dPr", f"{_M}naryPr", f"{_M}mPr",
    f"{_M}eqArrPr", f"{_M}funcPr", f"{_M}limLowPr", f"{_M}limUppPr",
    f"{_M}accPr", f"{_M}barPr", f"{_M}groupChrPr", f"{_M}ctrlPr",
    f"{_M}phantPr",
}


def _mval(el, default: str = "") -> str:
    """Return the m:val attribute value of an OMML element (namespace-qualified)."""
    if el is None:
        return default
    return el.get(f"{_M}val") or el.get("val") or default


def _mrow(inner: str) -> str:
    return f"<mrow>{inner}</mrow>"


def _children_ml(el) -> str:
    """Recursively convert all children of el to MathML string fragments."""
    return "" if el is None else "".join(_omml_el(c) for c in el)


def _get_math_variant(r_el) -> str:
    """
    Derive a MathML mathvariant string from an m:r element's run-property
    children.  Checks both m:rPr (OMML) and a:rPr (DrawingML).
    Returns '' when the MathML default (italic for single-char <mi>) applies.
    """
    _STY = {"p": "normal", "b": "bold", "bi": "bold-italic"}
    _SCR = {
        "cal": "script", "frak": "fraktur",
        "double-struck": "double-struck",
        "sans-serif": "sans-serif", "monospace": "monospace",
    }
    for child in r_el:
        if child.tag == f"{_M}rPr":
            if child.find(f"{_M}nor") is not None:
                return "normal"
            sty = child.find(f"{_M}sty")
            scr = child.find(f"{_M}scr")
            return _STY.get(_mval(sty), "") or _SCR.get(_mval(scr), "")
    return ""


def _math_token(text: str, variant: str) -> str:
    """Wrap a math text string in the appropriate MathML token element."""
    safe = escape_html(text)
    if re.fullmatch(r'\d+(?:\.\d*)?', text):
        return f"<mn>{safe}</mn>"
    if len(text) == 1 and text in _MATH_OPS:
        return f"<mo>{safe}</mo>"
    attr = f' mathvariant="{variant}"' if variant and variant != "italic" else ""
    return f"<mi{attr}>{safe}</mi>"


def _omml_el(el) -> str:
    """Convert a single OMML element (and subtree) to a MathML string."""
    tag = el.tag

    # Skip property / control elements entirely
    if tag in _M_PROP_TAGS:
        return ""

    # ── Math run (leaf) ──────────────────────────────────────────────────────
    if tag == f"{_M}r":
        t_el = el.find(f"{_M}t")
        text = (t_el.text or "") if t_el is not None else ""
        return _math_token(text, _get_math_variant(el)) if text else ""

    # ── Fraction ─────────────────────────────────────────────────────────────
    if tag == f"{_M}f":
        return f"<mfrac>{_mrow(_children_ml(el.find(f'{_M}num')))}{_mrow(_children_ml(el.find(f'{_M}den')))}</mfrac>"

    # ── Sub + Sup ─────────────────────────────────────────────────────────────
    if tag == f"{_M}sSubSup":
        b = _mrow(_children_ml(el.find(f"{_M}e")))
        s = _mrow(_children_ml(el.find(f"{_M}sub")))
        p = _mrow(_children_ml(el.find(f"{_M}sup")))
        return f"<msubsup>{b}{s}{p}</msubsup>"

    if tag == f"{_M}sSub":
        return f"<msub>{_mrow(_children_ml(el.find(f'{_M}e')))}{_mrow(_children_ml(el.find(f'{_M}sub')))}</msub>"

    if tag == f"{_M}sSup":
        return f"<msup>{_mrow(_children_ml(el.find(f'{_M}e')))}{_mrow(_children_ml(el.find(f'{_M}sup')))}</msup>"

    # ── Radical ───────────────────────────────────────────────────────────────
    if tag == f"{_M}rad":
        deg_str = _children_ml(el.find(f"{_M}deg"))
        base    = _mrow(_children_ml(el.find(f"{_M}e")))
        return f"<mroot>{base}{_mrow(deg_str)}</mroot>" if deg_str.strip() else f"<msqrt>{base}</msqrt>"

    # ── Delimiters / fence ────────────────────────────────────────────────────
    if tag == f"{_M}d":
        dPr = el.find(f"{_M}dPr")
        beg = escape_html(_mval(dPr.find(f"{_M}begChr") if dPr is not None else None, "("))
        end = escape_html(_mval(dPr.find(f"{_M}endChr") if dPr is not None else None, ")"))
        sep = escape_html(_mval(dPr.find(f"{_M}sepChr") if dPr is not None else None, "|"))
        inner_parts = [_mrow(_children_ml(e)) for e in el if e.tag == f"{_M}e"]
        inner = f"<mo>{sep}</mo>".join(inner_parts)
        return f"<mrow><mo>{beg}</mo>{inner}<mo>{end}</mo></mrow>"

    # ── N-ary operator (∫, ∑, ∏ …) ───────────────────────────────────────────
    if tag == f"{_M}nary":
        naryPr = el.find(f"{_M}naryPr")
        op_chr = escape_html(_mval(naryPr.find(f"{_M}chr") if naryPr is not None else None, "∫"))
        op_mo  = f'<mo largeop="true">{op_chr}</mo>'
        sub, sup = el.find(f"{_M}sub"), el.find(f"{_M}sup")
        body = _mrow(_children_ml(el.find(f"{_M}e")))
        if sub is not None and sup is not None:
            return f"<mrow><msubsup>{op_mo}{_mrow(_children_ml(sub))}{_mrow(_children_ml(sup))}</msubsup>{body}</mrow>"
        if sub is not None:
            return f"<mrow><msub>{op_mo}{_mrow(_children_ml(sub))}</msub>{body}</mrow>"
        if sup is not None:
            return f"<mrow><msup>{op_mo}{_mrow(_children_ml(sup))}</msup>{body}</mrow>"
        return f"<mrow>{op_mo}{body}</mrow>"

    # ── Function application ──────────────────────────────────────────────────
    if tag == f"{_M}func":
        return f"<mrow>{_children_ml(el.find(f'{_M}fName'))}<mo>&#x2061;</mo>{_children_ml(el.find(f'{_M}e'))}</mrow>"

    # ── Under / over limits ───────────────────────────────────────────────────
    if tag == f"{_M}limLow":
        return f"<munder>{_mrow(_children_ml(el.find(f'{_M}e')))}{_mrow(_children_ml(el.find(f'{_M}lim')))}</munder>"

    if tag == f"{_M}limUpp":
        return f"<mover>{_mrow(_children_ml(el.find(f'{_M}e')))}{_mrow(_children_ml(el.find(f'{_M}lim')))}</mover>"

    # ── Accent ────────────────────────────────────────────────────────────────
    if tag == f"{_M}acc":
        accPr = el.find(f"{_M}accPr")
        chr_v = escape_html(_mval(accPr.find(f"{_M}chr") if accPr is not None else None, "̂"))
        return f"<mover>{_mrow(_children_ml(el.find(f'{_M}e')))}<mo>{chr_v}</mo></mover>"

    # ── Bar ────────────────────────────────────────────────────────────────────
    if tag == f"{_M}bar":
        barPr = el.find(f"{_M}barPr")
        pos   = _mval(barPr.find(f"{_M}pos") if barPr is not None else None, "top")
        base  = _mrow(_children_ml(el.find(f"{_M}e")))
        bar   = '<mo stretchy="true">&#x305;</mo>'
        return f"<munder>{base}{bar}</munder>" if pos == "bot" else f"<mover>{base}{bar}</mover>"

    # ── Grouping character ────────────────────────────────────────────────────
    if tag == f"{_M}groupChr":
        gcPr  = el.find(f"{_M}groupChrPr")
        chr_v = escape_html(_mval(gcPr.find(f"{_M}chr") if gcPr is not None else None, "⏞"))
        pos   = _mval(gcPr.find(f"{_M}pos") if gcPr is not None else None, "top")
        base  = _mrow(_children_ml(el.find(f"{_M}e")))
        mo    = f"<mo>{chr_v}</mo>"
        return f"<munder>{base}{mo}</munder>" if pos == "bot" else f"<mover>{base}{mo}</mover>"

    # ── Matrix ────────────────────────────────────────────────────────────────
    if tag == f"{_M}m":
        rows = "".join(
            "<mtr>" + "".join(f"<mtd>{_mrow(_children_ml(e))}</mtd>" for e in mr if e.tag == f"{_M}e") + "</mtr>"
            for mr in el if mr.tag == f"{_M}mr"
        )
        return f"<mtable>{rows}</mtable>"

    # ── Equation array ────────────────────────────────────────────────────────
    if tag == f"{_M}eqArr":
        rows = "".join(f"<mtr><mtd>{_mrow(_children_ml(e))}</mtd></mtr>" for e in el if e.tag == f"{_M}e")
        return f"<mtable>{rows}</mtable>"

    # ── Phantom ───────────────────────────────────────────────────────────────
    if tag == f"{_M}phant":
        return f"<mphantom>{_mrow(_children_ml(el.find(f'{_M}e')))}</mphantom>"

    # ── Box / border box ──────────────────────────────────────────────────────
    if tag in (f"{_M}box", f"{_M}borderBox"):
        return f'<menclose notation="box">{_mrow(_children_ml(el.find(f"{_M}e")))}</menclose>'

    # ── Pre-scripts ───────────────────────────────────────────────────────────
    if tag == f"{_M}sPre":
        base = _children_ml(el.find(f"{_M}e"))
        sub  = _mrow(_children_ml(el.find(f"{_M}sub")))
        sup  = _mrow(_children_ml(el.find(f"{_M}sup")))
        return f"<mmultiscripts>{base}<mprescripts/>{sub}{sup}</mmultiscripts>"

    # ── oMath / oMathPara: transparent containers ─────────────────────────────
    if tag in (f"{_M}oMath", f"{_M}oMathPara"):
        return _children_ml(el)

    # ── Generic fallback: recurse (preserves content for unknown elements) ────
    return _children_ml(el)


def omml_to_mathml(omath_el, display: str = "inline", focusable: bool = False) -> str:
    """
    Convert an ``m:oMath`` (or ``m:oMathPara``) element to an HTML5-embeddable
    ``<math>`` string.

    Supports fractions, radicals, sub/sup, n-ary operators, delimiters,
    matrices, accents, limits, and more.  Unknown OMML elements fall back to
    processing their children so text content is always preserved.

    Args:
        omath_el:  ElementTree element whose tag is m:oMath or m:oMathPara.
        display:   ``"inline"`` or ``"block"``.
        focusable: When True, adds ``tabindex="0"`` so keyboard users can tab
                   to the equation and see a visible focus indicator.
    """
    inner = _children_ml(omath_el)
    tabindex = ' tabindex="0"' if focusable else ""
    return (
        f'<math xmlns="http://www.w3.org/1998/Math/MathML" display="{display}"{tabindex}>'
        f"<mrow>{inner}</mrow>"
        f"</math>"
    )


def _render_para_xml(para_p, bracket_refs: bool = True, focusable_math: bool = False) -> str:
    """
    Render an ``<a:p>`` element to HTML by walking its direct children.

    Handles interleaved text runs (``<a:r>``) and OMML math wrappers
    (``<a14:m>``).  Other children (paragraph properties, bookmarks, …)
    are silently skipped.
    """
    A_R   = f"{_A}r"
    A_T   = f"{_A}t"
    A_RPR = f"{_A}rPr"
    A14_M = f"{_A14}m"
    M_OMATH     = f"{_M}oMath"
    M_OMATHPARA = f"{_M}oMathPara"

    parts = []
    for child in para_p:
        tag = child.tag

        if tag == A_R:
            t_el = child.find(A_T)
            text = (t_el.text or "") if t_el is not None else ""
            if not text:
                continue
            safe = escape_html(text)
            rPr  = child.find(A_RPR)
            if rPr is not None:
                baseline = rPr.get("baseline")
                if baseline:
                    val = int(baseline)
                    if val > 0:
                        if bracket_refs and _CITATION_RE.match(text.strip()):
                            safe = f"<sup>[{safe}]</sup>"
                        else:
                            safe = f"<sup>{safe}</sup>"
                    elif val < 0:
                        safe = f"<sub>{safe}</sub>"
            parts.append(safe)

        elif tag == A14_M:
            # a14:m may contain m:oMathPara (block) or m:oMath (inline)
            omath_para = child.find(M_OMATHPARA)
            if omath_para is not None:
                for omath in omath_para:
                    if omath.tag == M_OMATH:
                        parts.append(omml_to_mathml(omath, display="block", focusable=focusable_math))
            else:
                omath = child.find(M_OMATH)
                if omath is not None:
                    parts.append(omml_to_mathml(omath, display="inline", focusable=focusable_math))

    return "".join(parts)


def render_txBody_from_xml(txBody_el, base_heading_level: int = 3, bracket_refs: bool = True, focusable_math: bool = False) -> str:
    """
    Render a raw ``<p:txBody>`` (or ``<a:txBody>``) element to HTML.

    Used for shapes extracted from ``<mc:AlternateContent>`` wrappers that
    are not accessible through the standard python-pptx shape iterator.
    """
    A_P   = f"{_A}p"
    A_PPR = f"{_A}pPr"
    A_BUAUTONUM = f"{_A}buAutoNum"
    A_BUCHAR    = f"{_A}buChar"
    A14_M = f"{_A14}m"

    html_parts = []
    list_tag   = ""

    for para_el in txBody_el.findall(A_P):
        # Text content (python-pptx-style: concatenate all t elements)
        text = "".join((t.text or "") for t in para_el.iter(f"{_A}t"))
        has_math = para_el.find(A14_M) is not None

        if not text.strip() and not has_math:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            continue

        safe = _render_para_xml(para_el, bracket_refs=bracket_refs, focusable_math=focusable_math)

        # Detect list type from pPr
        para_list_type = ""
        pPr = para_el.find(A_PPR)
        if pPr is not None:
            if pPr.find(A_BUAUTONUM) is not None:
                para_list_type = "ol"
            elif pPr.find(A_BUCHAR) is not None:
                para_list_type = "ul"

        if para_list_type:
            if list_tag and list_tag != para_list_type:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            if not list_tag:
                html_parts.append(f"<{para_list_type}>")
                list_tag = para_list_type
            html_parts.append(f"  <li>{safe}</li>")
        else:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            html_parts.append(f"<p>{safe}</p>")

    if list_tag:
        html_parts.append(f"</{list_tag}>")

    return "\n".join(html_parts)


def render_paragraph_runs(para, bracket_refs: bool = True, focusable_math: bool = False) -> str:
    """Render a paragraph's runs to HTML, preserving superscript/subscript
    and converting any embedded OMML math to MathML.

    If bracket_refs is True, superscript runs that look like citation numbers
    (digits, commas, spaces, hyphens) are wrapped in brackets: [1,2,3].
    """
    result = _render_para_xml(para._p, bracket_refs=bracket_refs, focusable_math=focusable_math)
    # Fallback for edge-case paragraphs with text but no structured runs/math
    if not result and para.text.strip():
        return escape_html(para.text.strip())
    return result


def _detect_list_type(para) -> str:
    """Detect paragraph list type: 'ol', 'ul', or '' (not a list item)."""
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        if para._p.pPr is not None:
            if para._p.pPr.find(f"{A_NS}buAutoNum") is not None:
                return "ol"
            if para._p.pPr.find(f"{A_NS}buChar") is not None:
                return "ul"
    except Exception:
        pass
    return ""


def render_text_frame(tf, base_heading_level: int = 3, bracket_refs: bool = True, focusable_math: bool = False) -> str:
    """
    Render a text frame to HTML.
    Paragraphs that look like headings (bold, larger font) become <hN>.
    Everything else becomes <p>.
    """
    html_parts = []
    list_tag = ""  # "", "ul", or "ol"

    for para in tf.paragraphs:
        text = para.text.strip()
        has_math = para._p.find(f"{_A14}m") is not None
        if not text and not has_math:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            continue

        safe = render_paragraph_runs(para, bracket_refs=bracket_refs, focusable_math=focusable_math)

        # Detect bullet / numbered list paragraph
        para_list_type = _detect_list_type(para)

        # Detect heading-like paragraph (bold + larger font in first run)
        is_heading = False
        try:
            if para.runs:
                first_run = para.runs[0]
                font_size = first_run.font.size
                is_bold = first_run.font.bold
                if is_bold and font_size and font_size >= Pt(14):
                    is_heading = True
        except Exception:
            pass

        if is_heading:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            level = min(base_heading_level, 6)
            html_parts.append(f"<h{level}>{safe}</h{level}>")
        elif para_list_type:
            # Switch list type if changing between ul and ol
            if list_tag and list_tag != para_list_type:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            if not list_tag:
                html_parts.append(f"<{para_list_type}>")
                list_tag = para_list_type
            html_parts.append(f"  <li>{safe}</li>")
        else:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            html_parts.append(f"<p>{safe}</p>")

    if list_tag:
        html_parts.append(f"</{list_tag}>")

    return "\n".join(html_parts)


def render_table(shape) -> str:
    """Render a table shape to an accessible HTML <table>."""
    table = shape.table
    rows_html = []
    for row_idx, row in enumerate(table.rows):
        cells_html = []
        for cell in row.cells:
            text = escape_html(cell.text.strip())
            tag = "th" if row_idx == 0 else "td"
            scope = ' scope="col"' if row_idx == 0 else ""
            cells_html.append(f"<{tag}{scope}>{text}</{tag}>")
        rows_html.append("<tr>" + "".join(cells_html) + "</tr>")

    caption = ""
    try:
        if shape.name:
            caption = f"<caption>{escape_html(shape.name)}</caption>"
    except Exception:
        pass

    return (
        '<table>\n'
        f'  {caption}\n'
        '  <thead>\n'
        f'    {rows_html[0]}\n'
        '  </thead>\n'
        '  <tbody>\n'
        + "\n".join(f"    {r}" for r in rows_html[1:])
        + "\n  </tbody>\n</table>"
    )


def merge_adjacent_lists(html: str) -> str:
    """Merge consecutive same-type list closings/openings with only whitespace between."""
    html = re.sub(r'</ul>\s*<ul>', '', html)
    html = re.sub(r'</ol>\s*<ol>', '', html)
    return html


# ── Core conversion ───────────────────────────────────────────────────────────

def get_slide_dimensions(prs):
    """Return slide canvas width and height in EMUs."""
    return prs.slide_width, prs.slide_height


def convert_slide(slide, slide_number: int, include_notes: bool, slide_width=None, slide_height=None, deck_title: str = "", bracket_refs: bool = True, focusable_math: bool = False, img_scale: float = 1.0) -> tuple:
    """Convert a single slide, returning (title_text, content_html_parts).

    Uses a single ordered walk over the slide's spTree so that
    <mc:AlternateContent> math shapes (invisible to python-pptx's slide.shapes)
    are emitted in document order alongside regular shapes rather than being
    appended at the end.

    The caller is responsible for section wrapping and title deduplication.
    """
    parts = []
    title_text = ""

    # ── Step 1: build shape_id → Shape map from the python-pptx API ──
    # Also extract the title and mark shapes to suppress (slide numbers,
    # running headers) by simply omitting them from the map.
    shape_map: dict = {}
    for shape in slide.shapes:
        if is_title_placeholder(shape) and shape_has_text(shape):
            title_text = shape.text_frame.text.strip()
        elif is_slide_number_placeholder(shape):
            pass  # skip — not added to map
        elif is_running_header(shape, deck_title):
            pass  # skip — not added to map
        else:
            shape_map[shape.shape_id] = shape

    # ── Step 2: helper to read the shape id from a raw XML element ──
    def _xml_shape_id(el) -> int:
        """Return the integer shape id from any spTree child element, or -1."""
        for nvpr_tag in (
            f"{_P}nvSpPr",
            f"{_P}nvPicPr",
            f"{_P}nvGrpSpPr",
            f"{_P}nvGraphicFramePr",
            f"{_P}nvCxnSpPr",
        ):
            nvpr = el.find(nvpr_tag)
            if nvpr is not None:
                cnvpr = nvpr.find(f"{_P}cNvPr")
                if cnvpr is not None:
                    try:
                        return int(cnvpr.get("id", -1))
                    except (ValueError, TypeError):
                        pass
        return -1

    # ── Step 3: helper to render one Shape object (reused for regular shapes) ──
    def _render_shape(shape) -> None:
        # ── Image ──
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            alt, is_decorative = get_alt_text_and_decorative(shape)
            if is_decorative:
                parts.append(f'  <!-- Decorative image skipped (slide {slide_number}) -->')
                return
            if not alt:
                alt = f"Image on slide {slide_number}"
            try:
                img = shape.image
                data_uri = image_to_data_uri(img.blob, img.content_type)
                px_w = round(shape.width / 914400 * 96 * img_scale)
                px_h = round(shape.height / 914400 * 96 * img_scale)
                size_style = f' style="width:{px_w}px; height:{px_h}px; max-width:100%; height:auto;"'
                parts.append(
                    f'  <figure>\n'
                    f'    <img src="{data_uri}" alt="{escape_html(alt)}"{size_style}>\n'
                    f'  </figure>'
                )
            except Exception as e:
                parts.append(f'  <!-- Image could not be embedded: {e} -->')

        # ── Table ──
        elif shape.has_table:
            parts.append(render_table(shape))

        # ── Code block (monospace font or shape named "code") ──
        elif shape.has_text_frame and is_code_shape(shape):
            rendered = render_code_block(shape.text_frame)
            if rendered:
                parts.append(rendered)

        # ── Text frame ──
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                return
            if is_subtitle_placeholder(shape):
                parts.append(f"  <p class=\"subtitle\">{escape_html(text)}</p>")
            else:
                rendered = render_text_frame(shape.text_frame, base_heading_level=3, bracket_refs=bracket_refs, focusable_math=focusable_math)
                if rendered:
                    parts.append(rendered)

        # ── Group shape — recurse for nested images ──
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    alt, is_decorative = get_alt_text_and_decorative(child)
                    if is_decorative:
                        continue
                    if not alt:
                        alt = f"Image on slide {slide_number}"
                    try:
                        img = child.image
                        data_uri = image_to_data_uri(img.blob, img.content_type)
                        size_style = ""
                        if slide_width and slide_height and slide_width > 0 and slide_height > 0:
                            px_w = round(child.width / 914400 * 96 * img_scale)
                            px_h = round(child.height / 914400 * 96 * img_scale)
                            size_style = f' style="width:{px_w}px; height:{px_h}px; max-width:100%; height:auto;"'
                        parts.append(
                            f'  <figure>\n'
                            f'    <img src="{data_uri}" alt="{escape_html(alt)}"{size_style}>\n'
                            f'  </figure>'
                        )
                    except Exception:
                        pass

    # ── Step 4: walk spTree direct children in document order ──
    # Regular shape tags dispatch via shape_map to _render_shape (keeps the
    # full python-pptx API).  <mc:AlternateContent> wrappers (math shapes
    # invisible to slide.shapes) are rendered inline via the XML renderer,
    # preserving their position in the reading order.
    _SHAPE_TAGS = {
        f"{_P}sp",
        f"{_P}pic",
        f"{_P}grpSp",
        f"{_P}graphicFrame",
        f"{_P}cxnSp",
    }
    try:
        spTree = slide._element.find(f".//{_P}spTree")
        if spTree is not None:
            for child in spTree:
                tag = child.tag

                if tag in _SHAPE_TAGS:
                    shape_id = _xml_shape_id(child)
                    shape = shape_map.get(shape_id)
                    if shape is not None:
                        _render_shape(shape)

                elif tag == f"{_MC}AlternateContent":
                    # Math-containing text box — render from XML in document order
                    choice = child.find(f"{_MC}Choice")
                    if choice is None:
                        continue
                    for sp_el in choice.findall(f"{_P}sp"):
                        txBody_el = sp_el.find(f"{_P}txBody")
                        if txBody_el is None:
                            continue
                        # Skip title placeholders (capture text if not yet found)
                        ph_type = ""
                        nvPr = sp_el.find(f".//{_P}nvPr")
                        if nvPr is not None:
                            ph = nvPr.find(f"{_P}ph")
                            if ph is not None:
                                ph_type = ph.get("type", "")
                        if ph_type in ("title", "ctrTitle"):
                            if not title_text:
                                title_text = "".join(
                                    (t.text or "")
                                    for t in txBody_el.iter(f"{_A}t")
                                ).strip()
                            continue
                        rendered = render_txBody_from_xml(txBody_el, bracket_refs=bracket_refs, focusable_math=focusable_math)
                        if rendered:
                            parts.append(rendered)
                # All other tags (p:grpSpPr etc.) are silently ignored
    except Exception:
        pass

    # ── Speaker notes ──
    if include_notes:
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(
                    f'  <aside class="speaker-notes">\n'
                    f'    <h3>Speaker Notes</h3>\n'
                    f'    <p>{escape_html(notes_text)}</p>\n'
                    f'  </aside>'
                )
        except Exception:
            pass

    return (title_text, parts)


def convert_pptx(input_path: Path, output_path: Path, include_notes: bool, bracket_refs: bool = True, focusable_math: bool = False, mathjax: bool = False, img_scale: float = 1.0) -> None:
    """Convert a .pptx file to an accessible HTML file."""
    prs = Presentation(str(input_path))
    title = input_path.stem.replace("_", " ").replace("-", " ").title()

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Extract deck title from slide 1 to filter running headers on other slides
    deck_title = ""
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if is_title_placeholder(shape) and shape_has_text(shape):
                deck_title = shape.text_frame.text.strip()
                break

    # Collect (title, content_parts, slide_number) for each slide
    slide_data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_title, parts = convert_slide(slide, i, include_notes, slide_width, slide_height, deck_title=deck_title, bracket_refs=bracket_refs, focusable_math=focusable_math, img_scale=img_scale)
        slide_data.append((slide_title, parts, i))

    # Group consecutive slides that share the same title
    # Each group becomes one <section> with one <h2>
    groups = []  # list of (title, first_slide_number, combined_parts)
    for slide_title, parts, slide_num in slide_data:
        effective_title = slide_title if slide_title else f"Slide {slide_num}"
        if groups and groups[-1][0] == effective_title:
            # Same title as previous — append content to existing group
            groups[-1][2].extend(parts)
        else:
            groups.append((effective_title, slide_num, list(parts)))

    # Build section HTML for each group, merging adjacent lists
    slide_sections = []
    for group_title, first_slide_num, content_parts in groups:
        section_lines = [f'<section aria-label="Slide {first_slide_num}">']
        section_lines.append(f"  <h2>{escape_html(group_title)}</h2>")
        content_html = "\n".join(content_parts)
        content_html = merge_adjacent_lists(content_html)
        section_lines.append(content_html)
        section_lines.append("</section>")
        slide_sections.append("\n".join(section_lines))

    slides_html = "\n\n".join(slide_sections)
    total = len(prs.slides)

    # ── Optional CSS / script injections ──────────────────────────────────────
    _focusable_css = (
        "\n    /* ── Focusable math ── */\n"
        '    math[tabindex="0"] { cursor: default; border-radius: 2px; outline-offset: 3px; }\n'
        '    math[tabindex="0"]:focus { outline: 2px solid #005fcc; background: #f0f4ff; }'
    ) if focusable_math else ""

    _mathjax_scripts = (
        '\n<script src="https://cdn.jsdelivr.net/npm/mathjax@4/tex-mml-chtml.js"'
        " defer></script>"
    ) if mathjax else ""

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{escape_html(title)}</title>
  <style>
    /* ── Base ── */
    *, *::before, *::after {{ box-sizing: border-box; }}
    body {{
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
      font-size: 1rem;
      line-height: 1.6;
      color: #1a1a1a;
      background: #f5f5f5;
      margin: 0;
      padding: 2rem 1rem;
    }}
    /* ── Skip link ── */
    .skip-link {{
      position: absolute;
      top: -40px;
      left: 0;
      background: #005fcc;
      color: #fff;
      padding: 0.5rem 1rem;
      border-radius: 0 0 4px 0;
      text-decoration: none;
      font-weight: bold;
      z-index: 100;
    }}
    .skip-link:focus {{ top: 0; }}
    /* ── Layout ── */
    main {{ max-width: 900px; margin: 0 auto; }}
    h1 {{
      font-size: 2rem;
      margin-bottom: 0.25rem;
      color: #111;
    }}
    .deck-meta {{
      color: #555;
      margin-bottom: 2rem;
      font-size: 0.95rem;
    }}
    /* ── Slide sections ── */
    section {{
      background: #fff;
      border: 1px solid #ddd;
      border-radius: 8px;
      padding: 2rem;
      margin-bottom: 2rem;
      box-shadow: 0 1px 4px rgba(0,0,0,.06);
    }}
    section h2 {{
      font-size: 1.5rem;
      margin-top: 0;
      color: #003366;
      border-bottom: 2px solid #e0e0e0;
      padding-bottom: 0.5rem;
    }}
    section h3 {{
      font-size: 1.15rem;
      color: #1a1a1a;
    }}
    p {{ margin: 0.5rem 0; }}
    ul {{ padding-left: 1.5rem; margin: 0.5rem 0; }}
    li {{ margin: 0.25rem 0; }}
    /* ── Code blocks ── */
    pre {{
      background: #1e1e1e;
      color: #d4d4d4;
      border-radius: 6px;
      padding: 1rem 1.25rem;
      overflow-x: auto;
      margin: 1rem 0;
      font-size: 0.9rem;
      line-height: 1.5;
    }}
    pre code {{
      font-family: Consolas, "Courier New", "Fira Mono", monospace;
      white-space: pre;
    }}
    /* ── Images ── */
    figure {{
      margin: 1rem 0;
      text-align: center;
    }}
    figure img {{
      max-width: 100%;
      height: auto;
      border-radius: 4px;
      border: 1px solid #e0e0e0;
    }}
    /* ── Tables ── */
    table {{
      border-collapse: collapse;
      width: 100%;
      margin: 1rem 0;
      font-size: 0.95rem;
    }}
    caption {{
      font-weight: bold;
      margin-bottom: 0.5rem;
      text-align: left;
    }}
    th, td {{
      border: 1px solid #ccc;
      padding: 0.5rem 0.75rem;
      text-align: left;
    }}
    th {{
      background: #003366;
      color: #fff;
    }}
    tr:nth-child(even) td {{ background: #f9f9f9; }}
    /* ── Speaker notes ── */
    .speaker-notes {{
      margin-top: 1.5rem;
      padding: 1rem;
      background: #fffbea;
      border-left: 4px solid #f0c040;
      border-radius: 4px;
    }}
    .speaker-notes h3 {{
      margin-top: 0;
      font-size: 0.9rem;
      text-transform: uppercase;
      letter-spacing: 0.05em;
      color: #7a6000;
    }}
    /* ── Subtitle ── */
    p.subtitle {{
      font-size: 1.1rem;
      color: #444;
      font-style: italic;
    }}{_focusable_css}
  </style>
{_mathjax_scripts}</head>
<body>
  <a href="#main-content" class="skip-link">Skip to main content</a>
  <main id="main-content">
    <h1>{escape_html(title)}</h1>
    <p class="deck-meta">Converted from <strong>{escape_html(input_path.name)}</strong> — {total} slide{"s" if total != 1 else ""}</p>

{slides_html}

  </main>
</body>
</html>
"""

    output_path.write_text(html, encoding="utf-8")
    print(f"OK  {input_path.name}  ->  {output_path}")


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint (.pptx) to accessible HTML, preserving alt text."
    )
    parser.add_argument(
        "input",
        help="Path to a .pptx file, or a folder containing .pptx files."
    )
    parser.add_argument(
        "-o", "--output",
        help="Output HTML file path (single-file mode only).",
        default=None
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes in the HTML output."
    )
    parser.add_argument(
        "--no-bracket-refs",
        action="store_true",
        help="Disable bracketing of superscript citation numbers (e.g. [1,2,3])."
    )
    parser.add_argument(
        "--focusable-math",
        action="store_true",
        help="Add tabindex='0' to each <math> element so keyboard users can tab to equations."
    )
    parser.add_argument(
        "--mathjax",
        action="store_true",
        help="Inject MathJax 4 (mml-chtml) script tags for enhanced math rendering in all browsers."
    )
    parser.add_argument(
        "--img-scale",
        type=float,
        default=100.0,
        metavar="PCT",
        help="Scale images to PCT%% of their slide dimensions (default: 100). "
             "Use values like 50 or 75 to reduce image size."
    )
    args = parser.parse_args()

    input_path    = Path(args.input)
    bracket_refs  = not args.no_bracket_refs
    focusable_math = args.focusable_math
    mathjax        = args.mathjax
    img_scale      = max(0.01, args.img_scale) / 100.0

    # ── Batch mode: folder ──
    if input_path.is_dir():
        pptx_files = sorted(input_path.glob("*.pptx"))
        if not pptx_files:
            print(f"No .pptx files found in: {input_path}")
            sys.exit(1)
        for pptx_file in pptx_files:
            out = pptx_file.with_suffix(".html")
            convert_pptx(pptx_file, out, args.include_notes, bracket_refs=bracket_refs, focusable_math=focusable_math, mathjax=mathjax, img_scale=img_scale)
        print(f"\nDone. {len(pptx_files)} file(s) converted.")
        return

    # ── Single-file mode ──
    if not input_path.exists():
        print(f"ERROR: File not found: {input_path}")
        sys.exit(1)
    if input_path.suffix.lower() != ".pptx":
        print(f"ERROR: Expected a .pptx file, got: {input_path.suffix}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_suffix(".html")

    convert_pptx(input_path, output_path, args.include_notes, bracket_refs=bracket_refs, focusable_math=focusable_math, mathjax=mathjax, img_scale=img_scale)
    print("Done.")


if __name__ == "__main__":
    main()
